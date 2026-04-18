from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── palette ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
ACCENT  = RGBColor(0xFF, 0x6B, 0x35)   # warm orange
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)   # teal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xB2, 0xC8)
YELLOW  = RGBColor(0xFF, 0xD6, 0x6B)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    # full-bleed background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def add_para(tf, text, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p

def accent_bar(slide, left=0.5, top=1.3, w=0.06, h=0.55, color=ACCENT):
    add_rect(slide, left, top, w, h, color)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
# top stripe
add_rect(s, 0, 0, 13.33, 0.12, ACCENT)
# bottom stripe
add_rect(s, 0, 7.38, 13.33, 0.12, ACCENT2)

add_text(s, "HACKATHON  ELECTRON  2026",
         0.7, 0.35, 12, 0.6,
         font_size=13, color=GRAY, align=PP_ALIGN.LEFT)

add_text(s, "Aventura noastră în\nAI & Machine Learning",
         0.7, 0.9, 11, 2.4,
         font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(s, 0.7, 3.5, 2.2, 0.06, ACCENT)

add_text(s, "Prezentare competitori  ·  ~4 minute",
         0.7, 3.65, 8, 0.5,
         font_size=16, color=GRAY, italic=True)

# decorative circles
for cx, cy, cr, col in [
    (11.5, 5.5, 1.8, RGBColor(0x1A, 0x2A, 0x4A)),
    (12.2, 4.2, 1.1, RGBColor(0x16, 0x24, 0x3E)),
    (10.3, 6.3, 0.9, RGBColor(0xFF,0x6B,0x35)),
]:
    sh = s.shapes.add_shape(9, Inches(cx-cr/2), Inches(cy-cr/2),
                             Inches(cr), Inches(cr))
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background()

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — De ce Electron?
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT)
accent_bar(s, 0.55, 0.3, 0.06, 0.6)

add_text(s, "De ce Electron?", 0.75, 0.25, 10, 0.9,
         font_size=36, bold=True)

add_rect(s, 0.75, 1.3, 11.5, 0.04, RGBColor(0x2A,0x3A,0x5A))

txBox = s.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.5), Inches(5.5))
tf = txBox.text_frame; tf.word_wrap = True

add_para(tf, "🎯  Curiozitate mai mare decât frica de eșec",
         font_size=22, bold=True, color=ACCENT)
add_para(tf, "Am vazut 'AI & ML' in titlu si pur si simplu n-am putut sa dam pass.",
         font_size=18, color=GRAY, space_before=2)

add_para(tf, "", font_size=10)

add_para(tf, "⚡  Provocarea concretă",
         font_size=22, bold=True, color=ACCENT2)
add_para(tf, "Probleme cu rețele neurale, clasificatoare de imagini, atacuri adversariale —\n"
             "exact genul de lucruri pe care le citești în papers și vrei să le atingi cu mâna.",
         font_size=18, color=GRAY, space_before=2)

add_para(tf, "", font_size=10)

add_para(tf, "🕐  Managementul timpului ca skill real",
         font_size=22, bold=True, color=YELLOW)
add_para(tf, "Știam că va fi intens. Voiam să vedem cât de bine ne descurcăm sub presiune.",
         font_size=18, color=GRAY, space_before=2)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Cu ce ne-am confruntat
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT2)
add_text(s, "Cu ce ne-am confruntat", 0.55, 0.25, 12, 0.9,
         font_size=36, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

# 3 cards
card_data = [
    (ACCENT,  "Adversarial Patches",
     "Construiește un patch 32×32 care,\nlipit pe o imagine, păcălește un\nclasificator să predicte clasa 42\ncu un confidence exact."),
    (ACCENT2, "Black-box Constraints",
     "Model weights disponibile, dar\ncondiția de confidence era secretă.\nTrebuia ghicit intervalul corect\nprin trial & error sistematic."),
    (YELLOW,  "Timp limitat",
     "Fiecare submit îți dă un răspuns\nbinar: da/nu. Zero detalii extra.\nIterezi repede sau pierzi timp."),
]
for i, (col, title, body) in enumerate(card_data):
    x = 0.55 + i * 4.2
    add_rect(s, x, 1.5, 3.9, 5.5, RGBColor(0x16,0x22,0x3A))
    add_rect(s, x, 1.5, 3.9, 0.1, col)
    add_text(s, title, x+0.15, 1.7, 3.6, 0.7,
             font_size=20, bold=True, color=col)
    add_text(s, body, x+0.15, 2.55, 3.6, 4.0,
             font_size=16, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Procesul de gândire
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, YELLOW)
add_text(s, "Cum ne-am gândit la probleme", 0.55, 0.25, 12, 0.9,
         font_size=34, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

steps = [
    ("01", ACCENT,  "Citim cerința de 3 ori",
     "Prima dată înțelegem ce vrea. A doua oară ce nu spune. A treia — ce capcane are."),
    ("02", ACCENT2, "Împărțim problema",
     "Input → model → output. Unde putem interveni? Ce este fix, ce este variabil?"),
    ("03", YELLOW,  "Prototip rapid, submit rapid",
     "Un patch random ca baseline. Dacă trece: noroc. Dacă nu: avem un punct de plecare."),
    ("04", WHITE,   "Iterăm pe date",
     "Ajustăm parametri, monitorizăm răspunsul verifierului, îngustăm intervalul de căutare."),
]

for i, (num, col, title, desc) in enumerate(steps):
    y = 1.5 + i * 1.35
    add_rect(s, 0.55, y, 0.7, 0.9, col)
    add_text(s, num, 0.55, y+0.05, 0.7, 0.8,
             font_size=26, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.4, y, 4.5, 0.5,
             font_size=19, bold=True, color=col)
    add_text(s, desc, 1.4, y+0.45, 11, 0.75,
             font_size=15, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Ce a funcționat, ce nu
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT)
add_text(s, "Ce a funcționat  ·  Ce nu", 0.55, 0.25, 12, 0.9,
         font_size=34, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

# left panel — worked
add_rect(s, 0.55, 1.5, 5.8, 5.6, RGBColor(0x12,0x2A,0x22))
add_rect(s, 0.55, 1.5, 5.8, 0.1, ACCENT2)
add_text(s, "✓  A funcționat", 0.7, 1.6, 5.5, 0.6,
         font_size=20, bold=True, color=ACCENT2)

worked = [
    "Gradient-based optimization pe patch",
    "Binary search pe confidence threshold",
    "Separarea sarcinilor în echipă (model vs. infra)",
    "Submit-uri mici și frecvente ca feedback loop",
]
txBox = s.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.5))
tf = txBox.text_frame; tf.word_wrap = True
for item in worked:
    add_para(tf, f"→  {item}", font_size=17, color=WHITE, space_before=8)

# right panel — didn't work
add_rect(s, 6.7, 1.5, 5.8, 5.6, RGBColor(0x2A,0x12,0x12))
add_rect(s, 6.7, 1.5, 5.8, 0.1, ACCENT)
add_text(s, "✗  N-a funcționat", 6.85, 1.6, 5.5, 0.6,
         font_size=20, bold=True, color=ACCENT)

failed = [
    "Patch-uri random fără direcție — pierdere de timp",
    "Over-engineering de la început",
    "Debugging fără logging → nu știam ce se întâmplă",
    "Subestimat cât durează setup-ul de mediu",
]
txBox = s.shapes.add_textbox(Inches(6.85), Inches(2.3), Inches(5.5), Inches(4.5))
tf = txBox.text_frame; tf.word_wrap = True
for item in failed:
    add_para(tf, f"→  {item}", font_size=17, color=GRAY, space_before=8)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Instrumente & Tehnologii
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT2)
add_text(s, "Instrumente & Tehnologii", 0.55, 0.25, 12, 0.9,
         font_size=34, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

tools = [
    (ACCENT,  "PyTorch",         "Backbone pentru orice — forward pass,\ngradienți, optimizare patch"),
    (ACCENT2, "Python + NumPy",  "Procesare imagini, manipulare tensori,\nscripturi rapide de test"),
    (YELLOW,  "PIL / OpenCV",    "Generare și inspecție vizuală a patch-urilor\nînainte de submit"),
    (WHITE,   "Git",             "Versionat fiecare iterație —\nnu pierzi ce a funcționat"),
    (ACCENT,  "VS Code",         "Live Share pentru colaborare\nîn timp real"),
    (ACCENT2, "Claude / ChatGPT","Explicat concepte noi rapid,\ndebug cu un al doilea ochi"),
]

for i, (col, name, desc) in enumerate(tools):
    col_idx = i % 3
    row_idx = i // 3
    x = 0.55 + col_idx * 4.2
    y = 1.5 + row_idx * 2.7
    add_rect(s, x, y, 3.9, 2.4, RGBColor(0x16,0x22,0x3A))
    add_rect(s, x, y, 0.1, 2.4, col)
    add_text(s, name, x+0.25, y+0.2, 3.5, 0.6,
             font_size=20, bold=True, color=col)
    add_text(s, desc, x+0.25, y+0.8, 3.5, 1.4,
             font_size=15, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI ne-a ajutat?
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, YELLOW)
add_text(s, "AI ne-a ajutat... sau nu?", 0.55, 0.25, 12, 0.9,
         font_size=34, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

add_text(s, "Răspunsul sincer: da și nu, depinde cum l-ai folosit.",
         0.55, 1.45, 12, 0.6, font_size=20, color=ACCENT, italic=True)

txBox = s.shapes.add_textbox(Inches(0.55), Inches(2.1), Inches(12), Inches(5))
tf = txBox.text_frame; tf.word_wrap = True

add_para(tf, "✅  Util când:", font_size=21, bold=True, color=ACCENT2)
add_para(tf, "→  Aveam nevoie să înțelegem rapid un concept nou (adversarial perturbations, FGSM, PGD)",
         font_size=17, color=WHITE, space_before=4)
add_para(tf, "→  Voiam să vedem rapid un schelet de cod ca punct de plecare",
         font_size=17, color=WHITE, space_before=4)
add_para(tf, "→  Debugging: explicat un error message cu context",
         font_size=17, color=WHITE, space_before=4)

add_para(tf, "", font_size=10)
add_para(tf, "❌  Limitări reale:", font_size=21, bold=True, color=ACCENT)
add_para(tf, "→  Nu cunoaște modelul tău specific — nu poate rezolva problema pentru tine",
         font_size=17, color=GRAY, space_before=4)
add_para(tf, "→  Codul generat trebuia adaptat / verificat întotdeauna",
         font_size=17, color=GRAY, space_before=4)
add_para(tf, "→  Uneori mai rapid să citești documentația direct",
         font_size=17, color=GRAY, space_before=4)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Ce am descoperit
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT)
add_text(s, "Ce am descoperit", 0.55, 0.25, 12, 0.9,
         font_size=34, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

discoveries = [
    (ACCENT,  "Adversarial ML",
     "Un patch mic de 32×32 poate complet păcăli un clasificator de 50 de clase.\n"
     "Modelele nu 'văd' cum vedem noi — asta e fascinant și îngrijorător în același timp."),
    (ACCENT2, "Security ↔ ML",
     "Atacurile adversariale sunt un domeniu de cercetare activ, direct legat de\n"
     "siguranța sistemelor AI din producție. Nu e doar teorie."),
    (YELLOW,  "Domeniul preferat",
     "Computer Vision + Security a fost preferata — faptul că poți 'ataca' rețele\n"
     "neurale cu imagini pare magie neagră la prima vedere."),
]

for i, (col, title, body) in enumerate(discoveries):
    y = 1.5 + i * 1.8
    add_rect(s, 0.55, y, 0.12, 1.4, col)
    add_text(s, title, 0.85, y+0.05, 5, 0.6,
             font_size=21, bold=True, color=col)
    add_text(s, body, 0.85, y+0.6, 11.8, 1.0,
             font_size=16, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Lecții de time management
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 0.18, 7.5, ACCENT2)
add_text(s, "Ce ne-a învățat hackathonul cu adevărat", 0.55, 0.25, 12, 0.9,
         font_size=32, bold=True)
add_rect(s, 0.55, 1.28, 12, 0.04, RGBColor(0x2A,0x3A,0x5A))

add_text(s, "(spoiler: nu e despre ML)",
         0.55, 1.35, 12, 0.5, font_size=17, color=GRAY, italic=True)

lessons = [
    ("🕐", ACCENT,  "Timpul e finit",      "Prima oră: explorezi. Ora 2+: construiești. Știi când să oprești rabbit holes."),
    ("🎯", ACCENT2, "Focus > Perfecțiune",  "Un submit mediocru acum bate un submit perfect mâine. Feedback rapid contează."),
    ("👥", YELLOW,  "Comunicarea salvează", "Dacă doi oameni lucrează pe același lucru fără să știe → dublu timp pierdut."),
    ("🔄", WHITE,   "Iterații mici",        "Schimbă un singur lucru odată. Altfel nu știi ce a rezolvat problema."),
]

for i, (icon, col, title, desc) in enumerate(lessons):
    col_idx = i % 2
    row_idx = i // 2
    x = 0.55 + col_idx * 6.3
    y = 2.05 + row_idx * 2.5
    add_rect(s, x, y, 5.9, 2.1, RGBColor(0x16,0x22,0x3A))
    add_text(s, icon, x+0.15, y+0.15, 0.8, 0.8, font_size=28)
    add_text(s, title, x+1.0, y+0.2, 4.7, 0.6,
             font_size=19, bold=True, color=col)
    add_text(s, desc, x+1.0, y+0.8, 4.7, 1.1,
             font_size=15, color=GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Closing
# ════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 0.12, ACCENT)
add_rect(s, 0, 7.38, 13.33, 0.12, ACCENT2)

add_text(s, "A fost greu. A fost bun.", 0.7, 1.2, 12, 1.4,
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 4.5, 2.9, 4.3, 0.06, ACCENT)

add_text(s,
         "Hackathonul Electron nu ne-a dat soluții — ne-a dat probleme bune.\n"
         "Și asta e exact ce trebuia.",
         0.7, 3.1, 12, 1.2,
         font_size=20, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "Mulțumim  🚀", 0.7, 5.5, 12, 0.9,
         font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "Electron Hackathon  ·  2026", 0.7, 6.5, 12, 0.5,
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────────────────
out = "/home/user/electron/Prezentare_Electron_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
